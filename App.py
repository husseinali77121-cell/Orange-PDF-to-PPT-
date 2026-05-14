import streamlit as st
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

# إعدادات الصفحة في Streamlit
st.set_page_config(page_title="محول PDF إلى PowerPoint", layout="centered")

def create_presentation(pdf_file):
    prs = Presentation()
    
    with pdfplumber.open(pdf_file) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if not text:
                continue
            
            # تقسيم النص إلى أسطر
            lines = text.split('\n')
            
            # إضافة شريحة جديدة (Layout 1: عنوان ومحتوى)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # تعيين العنوان (أول سطر في الصفحة)
            title_shape = slide.shapes.title
            title_shape.text = lines[0] if lines else "بدون عنوان"
            
            # إضافة بقية النص في مربع المحتوى
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            if len(lines) > 1:
                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 0
                    p.font.size = Pt(18)

    # حفظ الملف في ذاكرة مؤقتة (Buffer)
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# واجهة المستخدم
st.title("📄 تحويل PDF إلى عرض تقديمي PPTX")
st.write("قم برفع ملف PDF وسنقوم بتنظيمه في شرائح PowerPoint مرتبة.")

uploaded_file = st.file_uploader("اختر ملف PDF", type="pdf")

if uploaded_file is not None:
    with st.spinner('جاري التحليل والتحويل...'):
        try:
            # معالجة الملف
            result_ppt = create_presentation(uploaded_file)
            
            st.success("تم التحويل بنجاح!")
            
            # زر التحميل
            st.download_button(
                label="📥 تحميل ملف PowerPoint",
                data=result_ppt,
                file_name="converted_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"حدث خطأ أثناء المعالجة: {e}")

st.info("ملاحظة: هذا البرنامج يعمل بشكل أفضل مع ملفات الـ PDF التي تحتوي على نصوص قابلة للتحديد.")
