import os
import re
import fitz
import streamlit as st
import pdfplumber
import arabic_reshaper

from bidi.algorithm import get_display
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from PIL import Image

# =========================================
# PAGE CONFIG
# =========================================

st.set_page_config(
    page_title="PDF To PowerPoint AI",
    layout="wide"
)

# =========================================
# TEMPLATE CONFIG
# =========================================

TEMPLATES = {
    "Modern": "templates/modern.pptx",
    "Dark": "templates/dark.pptx",
    "Medical": "templates/medical.pptx",
    "Corporate": "templates/corporate.pptx"
}

# =========================================
# RTL SUPPORT
# =========================================

def rtl_text(text):
    reshaped = arabic_reshaper.reshape(text)
    return get_display(reshaped)

# =========================================
# SMART TEXT SPLITTER
# =========================================

def split_text(lines, max_lines=8):
    chunks = []

    for i in range(0, len(lines), max_lines):
        chunk = lines[i:i + max_lines]
        chunks.append(chunk)

    return chunks

# =========================================
# IMAGE EXTRACTION
# =========================================

def extract_images_from_pdf(pdf_path):
    pdf = fitz.open(stream=pdf_path.read(), filetype="pdf")

    image_paths = []

    if not os.path.exists("extracted_images"):
        os.makedirs("extracted_images")

    for page_index in range(len(pdf)):
        page = pdf[page_index]
        images = page.get_images(full=True)

        for img_index, img in enumerate(images):
            xref = img[0]

            base_image = pdf.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            image_path = f"extracted_images/page{page_index+1}_{img_index}.{image_ext}"

            with open(image_path, "wb") as f:
                f.write(image_bytes)

            image_paths.append(image_path)

    return image_paths

# =========================================
# CREATE PRESENTATION
# =========================================

def create_presentation(uploaded_pdf, template_path, add_images=True):

    prs = Presentation(template_path)

    image_paths = extract_images_from_pdf(uploaded_pdf)

    uploaded_pdf.seek(0)

    with pdfplumber.open(uploaded_pdf) as pdf:

        total_pages = len(pdf.pages)

        progress_bar = st.progress(0)

        image_counter = 0

        for page_num, page in enumerate(pdf.pages):

            text = page.extract_text()

            if not text:
                continue

            lines = [line.strip() for line in text.split('\n') if line.strip()]

            if not lines:
                continue

            title = lines[0]

            content_lines = lines[1:]

            split_chunks = split_text(content_lines, max_lines=7)

            for chunk in split_chunks:

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                # ====================================
                # TITLE
                # ====================================

                title_box = slide.shapes.title

                title_box.text = rtl_text(title)

                title_para = title_box.text_frame.paragraphs[0]

                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.RIGHT

                # ====================================
                # BODY
                # ====================================

                body_shape = slide.placeholders[1]

                tf = body_shape.text_frame

                tf.clear()

                for line in chunk:

                    p = tf.add_paragraph()

                    p.text = rtl_text("• " + line)

                    p.font.size = Pt(20)

                    p.alignment = PP_ALIGN.RIGHT

                # ====================================
                # IMAGE INSERTION
                # ====================================

                if add_images and image_counter < len(image_paths):

                    try:
                        slide.shapes.add_picture(
                            image_paths[image_counter],
                            Inches(6),
                            Inches(1.5),
                            width=Inches(3)
                        )

                        image_counter += 1

                    except:
                        pass

            progress = int((page_num + 1) / total_pages * 100)

            progress_bar.progress(progress)

    ppt_io = BytesIO()

    prs.save(ppt_io)

    ppt_io.seek(0)

    return ppt_io

# =========================================
# STREAMLIT UI
# =========================================

st.title("📄 PDF To PowerPoint AI")

st.write(
    """
    تحويل ملفات PDF إلى عروض PowerPoint احترافية
    مع دعم:
    
    ✅ Templates متعددة  
    ✅ اللغة العربية  
    ✅ استخراج الصور  
    ✅ تقسيم ذكي للشرائح  
    """
)

# =========================================
# SIDEBAR
# =========================================

st.sidebar.header("⚙️ الإعدادات")

selected_template = st.sidebar.selectbox(
    "اختر Template",
    list(TEMPLATES.keys())
)

presentation_type = st.sidebar.radio(
    "نوع العرض",
    [
        "Medical",
        "Business",
        "Education",
        "Research"
    ]
)

add_images = st.sidebar.checkbox(
    "استخراج الصور من PDF",
    value=True
)

# =========================================
# FILE UPLOAD
# =========================================

uploaded_file = st.file_uploader(
    "قم برفع ملف PDF",
    type="pdf"
)

# =========================================
# CONVERT BUTTON
# =========================================

if uploaded_file is not None:

    st.success("تم رفع الملف بنجاح")

    if st.button("🚀 تحويل إلى PowerPoint"):

        with st.spinner("جاري التحويل..."):

            try:

                ppt_file = create_presentation(
                    uploaded_file,
                    TEMPLATES[selected_template],
                    add_images
                )

                st.success("✅ تم إنشاء العرض بنجاح")

                st.download_button(
                    label="📥 تحميل PowerPoint",
                    data=ppt_file,
                    file_name="AI_Presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            except Exception as e:

                st.error(f"حدث خطأ: {e}")

# =========================================
# FOOTER
# =========================================

st.markdown("---")

st.caption("Developed with Streamlit + Python PPTX + PyMuPDF")
